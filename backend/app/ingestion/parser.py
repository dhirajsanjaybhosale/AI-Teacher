import io
import os
import re
from typing import Dict, Any, List, Optional

try:
    import pymupdf as fitz
except ImportError:
    try:
        import fitz
    except ImportError:
        fitz = None

# Optional docx and pptx libraries
try:
    import docx
    _DOCX_AVAILABLE = True
except ImportError:
    _DOCX_AVAILABLE = False

try:
    import pptx
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


class DocumentParser:
    """
    Universal Document Ingestion Parser.
    Supports PDF, DOCX, PPTX, and TXT documents.
    Preserves page/slide numbers, detects headings, chapters, sections, tables, formulas,
    and detects scanned/image-only PDFs with clear user warnings.
    """

    SUPPORTED_EXTENSIONS = [".pdf", ".docx", ".pptx", ".txt", ".md"]

    def __init__(self):
        pass

    def extract_text_from_bytes(self, doc_bytes: bytes, filename: str = "document.pdf") -> Dict[str, Any]:
        """
        Parses raw bytes of a document based on its file extension.
        """
        ext = os.path.splitext(filename)[1].lower()
        if ext == ".docx":
            return self._parse_docx(doc_bytes, filename)
        elif ext == ".pptx":
            return self._parse_pptx(doc_bytes, filename)
        elif ext in [".txt", ".md"]:
            return self._parse_txt(doc_bytes, filename)
        else:
            # Default to PDF
            return self._parse_pdf(doc_bytes, filename)

    def extract_text_from_file(self, file_path: str) -> Dict[str, Any]:
        """
        Parses a local file from disk.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document file not found at: {file_path}")
        filename = os.path.basename(file_path)
        with open(file_path, "rb") as f:
            data = f.read()
        return self.extract_text_from_bytes(data, filename)

    # ------------------------------------------------------------------
    # 1. PDF Parser (PyMuPDF)
    # ------------------------------------------------------------------
    def _parse_pdf(self, pdf_bytes: bytes, filename: str) -> Dict[str, Any]:
        if fitz is None:
            raise RuntimeError("PyMuPDF (fitz) is not installed.")

        try:
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        except Exception as e:
            return {
                "filename": filename,
                "title": filename,
                "author": "Unknown",
                "total_pages": 0,
                "total_words": 0,
                "full_text": "",
                "pages": [],
                "toc": [],
                "is_scanned": False,
                "error": f"Malformed PDF: {str(e)}"
            }

        pages_content = []
        full_text_list = []
        toc = doc.get_toc() or []
        total_extracted_words = 0
        headings_detected = []
        formulas_detected = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text").strip()
            
            # Detect headings and structure via text blocks
            page_headings = []
            page_tables = []
            page_formulas = []
            
            try:
                blocks = page.get_text("blocks")
                for b in blocks:
                    block_text = b[4].strip()
                    # Check for formula indicators (LaTeX or math equations)
                    if any(sym in block_text for sym in ["=", "+", "-", "*", "/", "∫", "∑", "λ", "Ω", "π", "Δ", "±"]) and len(block_text.split()) <= 12:
                        page_formulas.append(block_text)
                    # Check for potential headings (short lines with title/upper case)
                    elif len(block_text.split()) <= 8 and (block_text.isupper() or block_text.istitle() or re.match(r'^(?:Chapter|Section|\d+\.|\d+\.\d+)', block_text)):
                        page_headings.append(block_text)
            except Exception:
                pass

            word_count = len(text.split())
            total_extracted_words += word_count

            if text:
                pages_content.append({
                    "page_number": page_num + 1,
                    "text": text,
                    "word_count": word_count,
                    "headings": page_headings,
                    "formulas": page_formulas
                })
                full_text_list.append(text)
                headings_detected.extend(page_headings)
                formulas_detected.extend(page_formulas)

        full_text = "\n\n".join(full_text_list)
        metadata = doc.metadata or {}
        title = metadata.get("title") or filename.replace(".pdf", "").replace("_", " ").title()

        # Scanned PDF Detection: Multi-page document with virtually zero extracted text
        is_scanned = len(doc) > 0 and total_extracted_words < 10
        scan_warning = ""
        if is_scanned:
            scan_warning = "This appears to be a scanned document. OCR is required to read it."

        return {
            "filename": filename,
            "title": title,
            "author": metadata.get("author", "Unknown"),
            "total_pages": len(doc),
            "total_words": total_extracted_words,
            "full_text": full_text,
            "pages": pages_content,
            "toc": toc,
            "headings": headings_detected[:20],
            "formulas": formulas_detected[:10],
            "is_scanned": is_scanned,
            "scan_warning": scan_warning
        }

    # ------------------------------------------------------------------
    # 2. DOCX Parser (python-docx)
    # ------------------------------------------------------------------
    def _parse_docx(self, doc_bytes: bytes, filename: str) -> Dict[str, Any]:
        if not _DOCX_AVAILABLE:
            raise RuntimeError("python-docx is not installed.")

        file_stream = io.BytesIO(doc_bytes)
        doc = docx.Document(file_stream)
        
        pages_content = []
        full_text_list = []
        headings = []
        current_section = "Introduction"
        current_page_num = 1
        current_page_paras = []

        for p in doc.paragraphs:
            text = p.text.strip()
            if not text:
                continue

            # Detect headings
            style_name = p.style.name.lower() if p.style else ""
            if "heading" in style_name or re.match(r'^(?:Chapter|Section|\d+\.|\d+\.\d+)', text):
                headings.append(text)
                current_section = text

            current_page_paras.append(text)
            full_text_list.append(text)

            # Approximate page breakdown every 400 words
            if sum(len(x.split()) for x in current_page_paras) >= 400:
                p_text = "\n\n".join(current_page_paras)
                pages_content.append({
                    "page_number": current_page_num,
                    "text": p_text,
                    "word_count": len(p_text.split()),
                    "section": current_section
                })
                current_page_num += 1
                current_page_paras = []

        if current_page_paras:
            p_text = "\n\n".join(current_page_paras)
            pages_content.append({
                "page_number": current_page_num,
                "text": p_text,
                "word_count": len(p_text.split()),
                "section": current_section
            })

        # Also extract table text
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    full_text_list.append(f"[Table Row]: {row_text}")

        full_text = "\n\n".join(full_text_list)
        title = filename.replace(".docx", "").replace("_", " ").title()

        return {
            "filename": filename,
            "title": title,
            "author": "Document Author",
            "total_pages": max(1, len(pages_content)),
            "total_words": len(full_text.split()),
            "full_text": full_text,
            "pages": pages_content,
            "headings": headings[:20],
            "is_scanned": False,
            "scan_warning": ""
        }

    # ------------------------------------------------------------------
    # 3. PPTX Parser (python-pptx)
    # ------------------------------------------------------------------
    def _parse_pptx(self, pptx_bytes: bytes, filename: str) -> Dict[str, Any]:
        if not _PPTX_AVAILABLE:
            raise RuntimeError("python-pptx is not installed.")

        file_stream = io.BytesIO(pptx_bytes)
        prs = pptx.Presentation(file_stream)
        pages_content = []
        full_text_list = []
        slide_titles = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            slide_texts = []
            slide_title = f"Slide {slide_num}"

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                            if shape == slide.shapes[0] and not slide_titles:
                                slide_title = text

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[Speaker Notes]: {notes}")

            if slide_texts:
                slide_titles.append(slide_title)
                slide_combined = "\n".join(slide_texts)
                pages_content.append({
                    "page_number": slide_num,
                    "text": slide_combined,
                    "word_count": len(slide_combined.split()),
                    "slide_title": slide_title
                })
                full_text_list.append(f"--- Slide {slide_num}: {slide_title} ---\n{slide_combined}")

        full_text = "\n\n".join(full_text_list)
        title = slide_titles[0] if slide_titles else filename.replace(".pptx", "").replace("_", " ").title()

        return {
            "filename": filename,
            "title": title,
            "author": "Presentation Author",
            "total_pages": len(prs.slides),
            "total_words": len(full_text.split()),
            "full_text": full_text,
            "pages": pages_content,
            "headings": slide_titles[:20],
            "is_scanned": False,
            "scan_warning": ""
        }

    # ------------------------------------------------------------------
    # 4. TXT / Markdown Parser
    # ------------------------------------------------------------------
    def _parse_txt(self, txt_bytes: bytes, filename: str) -> Dict[str, Any]:
        # Try common encodings
        for enc in ["utf-8", "utf-16", "latin-1"]:
            try:
                raw_text = txt_bytes.decode(enc)
                break
            except Exception:
                continue
        else:
            raw_text = txt_bytes.decode("utf-8", errors="replace")

        clean_text = raw_text.strip()
        paragraphs = [p.strip() for p in clean_text.split("\n\n") if p.strip()]
        
        pages_content = []
        current_page = 1
        page_paras = []
        headings = []

        for p in paragraphs:
            if p.startswith("#") or re.match(r'^(?:Chapter|Section|\d+\.)', p):
                headings.append(p.lstrip("#").strip())

            page_paras.append(p)
            if sum(len(x.split()) for x in page_paras) >= 400:
                p_str = "\n\n".join(page_paras)
                pages_content.append({
                    "page_number": current_page,
                    "text": p_str,
                    "word_count": len(p_str.split())
                })
                current_page += 1
                page_paras = []

        if page_paras:
            p_str = "\n\n".join(page_paras)
            pages_content.append({
                "page_number": current_page,
                "text": p_str,
                "word_count": len(p_str.split())
            })

        title = headings[0] if headings else filename.replace(".txt", "").replace(".md", "").replace("_", " ").title()

        return {
            "filename": filename,
            "title": title,
            "author": "Text Document",
            "total_pages": max(1, len(pages_content)),
            "total_words": len(clean_text.split()),
            "full_text": clean_text,
            "pages": pages_content,
            "headings": headings[:20],
            "is_scanned": False,
            "scan_warning": ""
        }


# Backward-compatible alias for existing imports
PDFParser = DocumentParser

